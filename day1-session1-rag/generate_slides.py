from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

BLUE = RGBColor(0x1A, 0x56, 0xDB)
DARK_BLUE = RGBColor(0x0E, 0x3B, 0x8C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x6B, 0x72, 0x80)
DARK = RGBColor(0x1F, 0x29, 0x37)
LIGHT_BG = RGBColor(0xF8, 0xFA, 0xFC)
GREEN = RGBColor(0x05, 0x96, 0x69)
ORANGE = RGBColor(0xD9, 0x77, 0x06)
PURPLE = RGBColor(0x7C, 0x3A, 0xED)
RED = RGBColor(0xDC, 0x26, 0x26)
TEAL = RGBColor(0x0D, 0x94, 0x88)
SNOW_BLUE = RGBColor(0x29, 0xB5, 0xE8)


def add_bg(slide, color=LIGHT_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False, color=DARK, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    return txBox


def add_bullet_slide(slide, items, start_top, left=Inches(1.5), font_size=20, color=DARK):
    for i, item in enumerate(items):
        txBox = slide.shapes.add_textbox(left, start_top + Inches(i * 0.55), Inches(13), Inches(0.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color


# ========== SLIDE 1: Title ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BLUE)
add_shape(slide, Inches(0), Inches(0), Inches(16), Inches(0.15), SNOW_BLUE)
add_shape(slide, Inches(0), Inches(8.85), Inches(16), Inches(0.15), SNOW_BLUE)

add_text_box(slide, Inches(2), Inches(1.5), Inches(12), Inches(1.5),
             "Snowflake Workshop x Bank Nusantara", 42, True, WHITE, PP_ALIGN.CENTER)
add_text_box(slide, Inches(2), Inches(3.2), Inches(12), Inches(1),
             "Day 1 - Session 1", 32, True, SNOW_BLUE, PP_ALIGN.CENTER)
add_text_box(slide, Inches(2), Inches(4.5), Inches(12), Inches(1.2),
             "RAG Applications + Snowflake Intelligence", 36, True, WHITE, PP_ALIGN.CENTER)
add_text_box(slide, Inches(2), Inches(6.0), Inches(12), Inches(0.6),
             "Membangun AI Banking Assistant dengan Snowflake Cortex AI", 20, False, RGBColor(0xBF, 0xDB, 0xFE), PP_ALIGN.CENTER)
add_text_box(slide, Inches(2), Inches(7.5), Inches(12), Inches(0.5),
             "Powered by Snowflake Cortex AI  |  April 2026", 16, False, GRAY, PP_ALIGN.CENTER)

# ========== SLIDE 2: Agenda ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(16), Inches(1.2), DARK_BLUE)
add_text_box(slide, Inches(1), Inches(0.2), Inches(14), Inches(0.8), "Agenda Session 1", 32, True, WHITE, PP_ALIGN.LEFT)

items = [
    "Part 1    Environment Setup & Data Loading",
    "             - Create database, schema, stage",
    "             - Upload PDF & CSV files, COPY INTO tables",
    "",
    "Part 2    Unstructured Data Processing & Cortex Search",
    "             - PDF extraction, text chunking, search service",
    "",
    "Part 3    Structured Data & Semantic View (Cortex Analyst)",
    "             - Create Semantic View dengan metrics perbankan",
    "",
    "Part 4    Cortex Agent - AI Banking Assistant",
    "             - Multi-tool agent: Search + Analyst + Web + Email",
    "",
    "Part 5    Snowflake Intelligence - Publish Agent",
    "             - Deploy agent untuk business users",
]
add_bullet_slide(slide, items, Inches(1.5), font_size=18)

# ========== SLIDE 3: What We Will Build ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(16), Inches(1.2), DARK_BLUE)
add_text_box(slide, Inches(1), Inches(0.2), Inches(14), Inches(0.8), "Apa yang Akan Kita Bangun?", 32, True, WHITE)

add_text_box(slide, Inches(1), Inches(1.5), Inches(14), Inches(0.6),
             "AI Banking Assistant yang mampu:", 22, True, DARK)

features = [
    "1.  Menjawab pertanyaan berdasarkan dokumen SOP internal bank (6 dokumen PDF)",
    "2.  Menganalisis data operasional perbankan: NPL, DPK, LDR, transaksi",
    "3.  Melakukan web search untuk regulasi & informasi terkini",
    "4.  Memberikan jawaban: Fact-Based, Diagnostic, Recommendation + Chart",
    "5.  Mengirim email notifikasi ke stakeholder",
    "6.  Di-publish sebagai aplikasi untuk Relationship Manager, Risk Analyst, dll.",
]
add_bullet_slide(slide, features, Inches(2.3), font_size=19)

add_text_box(slide, Inches(1), Inches(6.5), Inches(14), Inches(1),
             "Semua dibangun di dalam Snowflake - tanpa infrastruktur ML terpisah, tanpa API LLM external",
             18, True, BLUE, PP_ALIGN.CENTER)

# ========== SLIDE 4: Architecture ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(16), Inches(1.2), DARK_BLUE)
add_text_box(slide, Inches(1), Inches(0.2), Inches(14), Inches(0.8), "Reference Architecture", 32, True, WHITE)

img_path = os.path.join(os.path.dirname(__file__), "images", "architecture_day1_session1.png")
if os.path.exists(img_path):
    slide.shapes.add_picture(img_path, Inches(1), Inches(1.4), Inches(14), Inches(7.2))
else:
    add_text_box(slide, Inches(3), Inches(4), Inches(10), Inches(1),
                 "[Architecture Diagram - see images/architecture_day1_session1.png]", 20, False, GRAY, PP_ALIGN.CENTER)

# ========== SLIDE 5: Data Sources ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(16), Inches(1.2), DARK_BLUE)
add_text_box(slide, Inches(1), Inches(0.2), Inches(14), Inches(0.8), "Part 1: Data Sources", 32, True, WHITE)

add_text_box(slide, Inches(1), Inches(1.5), Inches(6.5), Inches(0.5), "Unstructured Data (PDF)", 22, True, ORANGE)
pdf_items = [
    "  01. SOP Pembukaan Rekening",
    "  02. Panduan KYC (Know Your Customer)",
    "  03. Kebijakan Kredit UMKM",
    "  04. Panduan Anti Money Laundering",
    "  05. Prosedur Restrukturisasi Kredit",
    "  06. Panduan Perbankan Digital",
]
add_bullet_slide(slide, pdf_items, Inches(2.2), left=Inches(1), font_size=17)

add_text_box(slide, Inches(8.5), Inches(1.5), Inches(6.5), Inches(0.5), "Structured Data (CSV)", 22, True, GREEN)
csv_items = [
    "  DIM_CABANG       (67 cabang)",
    "  DIM_NASABAH      (10,000 nasabah)",
    "  FACT_TRANSAKSI   (100,000 transaksi)",
    "  FACT_KREDIT      (5,000 kredit)",
    "  FACT_SIMPANAN    (15,000 simpanan)",
]
add_bullet_slide(slide, csv_items, Inches(2.2), left=Inches(8.5), font_size=17)

add_text_box(slide, Inches(1), Inches(6.8), Inches(14), Inches(1),
             "Semua data menggunakan terminologi perbankan Indonesia yang realistis",
             18, True, GRAY, PP_ALIGN.CENTER)

# ========== SLIDE 6: Part 1 Steps ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(16), Inches(1.2), DARK_BLUE)
add_text_box(slide, Inches(1), Inches(0.2), Inches(14), Inches(0.8), "Part 1: Environment Setup & Data Loading", 32, True, WHITE)

steps = [
    "Step 1.1   CREATE DATABASE & SCHEMA",
    "              - BANK_NUSANTARA_DB dengan schema RAW_DATA dan ANALYTICS",
    "",
    "Step 1.2   CREATE STAGE",
    "              - STG_DOCUMENTS (untuk PDF) dan STG_CSV_DATA (untuk CSV)",
    "",
    "Step 1.3   UPLOAD FILES ke Stage",
    "              - Via Snowsight UI atau Snowflake CLI",
    "",
    "Step 1.4   CREATE TABLE & COPY INTO",
    "              - 5 tabel: DIM_CABANG, DIM_NASABAH, FACT_TRANSAKSI, FACT_KREDIT, FACT_SIMPANAN",
    "",
    "Step 1.5   VERIFIKASI DATA",
    "              - Row count dan preview data",
]
add_bullet_slide(slide, steps, Inches(1.5), font_size=18)

# ========== SLIDE 7: Part 2 - Cortex Search ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(16), Inches(1.2), DARK_BLUE)
add_text_box(slide, Inches(1), Inches(0.2), Inches(14), Inches(0.8), "Part 2: Unstructured Data & Cortex Search", 32, True, WHITE)

steps = [
    "Step 2.1   PARSE_DOCUMENT - Extract text dari PDF",
    "              SNOWFLAKE.CORTEX.PARSE_DOCUMENT(@stage, file, {'mode': 'LAYOUT'})",
    "",
    "Step 2.2   SPLIT_TEXT_RECURSIVE_CHARACTER - Chunking",
    "              Memecah dokumen menjadi chunks 1500 chars dengan 300 chars overlap",
    "",
    "Step 2.3   CREATE CORTEX SEARCH SERVICE",
    "              Hybrid search (semantic + keyword) pada chunks dokumen",
    "",
    "Step 2.4   TEST - Coba tanya tentang SOP",
    '              "Apa syarat pembukaan rekening untuk WNA?"',
    '              "Bagaimana prosedur KYC untuk nasabah berisiko tinggi?"',
    '              "Apa red flags transaksi money laundering?"',
]
add_bullet_slide(slide, steps, Inches(1.5), font_size=18)

# ========== SLIDE 8: Cortex Search Flow ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(16), Inches(1.2), DARK_BLUE)
add_text_box(slide, Inches(1), Inches(0.2), Inches(14), Inches(0.8), "Cortex Search: How It Works", 32, True, WHITE)

flow_items = [
    ("PDF Documents", ORANGE, 1.5),
    ("PARSE_DOCUMENT\n(Extract Text)", BLUE, 3.0),
    ("SPLIT_TEXT\n(Chunking)", BLUE, 4.5),
    ("Cortex Search\nService", GREEN, 6.0),
    ("User Query\n(Natural Language)", PURPLE, 7.5),
]
for i, (text, color, top) in enumerate(flow_items):
    shape = add_shape(slide, Inches(5.5), Inches(top), Inches(5), Inches(1.1), RGBColor(0xF3, 0xF4, 0xF6))
    shape.line.color.rgb = color
    shape.line.width = Pt(2)
    add_text_box(slide, Inches(5.8), Inches(top + 0.15), Inches(4.4), Inches(0.8), text, 18, True, color, PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(2.5), Inches(4), Inches(5),
             "Key Concepts:\n\n"
             "- Layout mode preserves\n  document structure\n\n"
             "- Chunk size: 1500 chars\n  Overlap: 300 chars\n\n"
             "- Hybrid search combines\n  semantic + keyword\n\n"
             "- Auto-refresh: 1 hour lag",
             16, False, DARK)

# ========== SLIDE 9: Part 3 - Semantic View ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(16), Inches(1.2), DARK_BLUE)
add_text_box(slide, Inches(1), Inches(0.2), Inches(14), Inches(0.8), "Part 3: Semantic View & Cortex Analyst", 32, True, WHITE)

add_text_box(slide, Inches(1), Inches(1.5), Inches(7), Inches(0.5), "Semantic View Components:", 22, True, GREEN)

sv_items = [
    "TABLES (5 tables with column descriptions & synonyms)",
    "  - NASABAH, CABANG, TRANSAKSI, KREDIT, SIMPANAN",
    "",
    "RELATIONSHIPS (Foreign Keys)",
    "  - Nasabah -> Cabang, Transaksi -> Nasabah, dll.",
    "",
    "METRICS (8 banking metrics)",
    "  - TOTAL_DPK, NPL_RATIO, LDR, TOTAL_OUTSTANDING_KREDIT",
    "  - JUMLAH_NASABAH_AKTIF, TOTAL_VOLUME_TRANSAKSI, dll.",
    "",
    "FILTERS (5 named filters)",
    "  - NASABAH_AKTIF, TRANSAKSI_BERHASIL, KREDIT_BERMASALAH, dll.",
]
add_bullet_slide(slide, sv_items, Inches(2.2), font_size=17)

add_text_box(slide, Inches(1), Inches(7.5), Inches(14), Inches(0.5),
             "Cortex Analyst menggunakan Semantic View untuk menghasilkan SQL dari natural language",
             17, True, BLUE, PP_ALIGN.CENTER)

# ========== SLIDE 10: Metrics ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(16), Inches(1.2), DARK_BLUE)
add_text_box(slide, Inches(1), Inches(0.2), Inches(14), Inches(0.8), "Banking Metrics dalam Semantic View", 32, True, WHITE)

metrics = [
    ("TOTAL_DPK", "Dana Pihak Ketiga - Total saldo simpanan aktif", "SUM(saldo) WHERE status = 'Aktif'"),
    ("NPL_RATIO", "Non-Performing Loan Ratio (target < 5%)", "Outstanding kol 3,4,5 / Total Outstanding * 100"),
    ("LDR", "Loan to Deposit Ratio (target 78-92%)", "Total Kredit / Total DPK * 100"),
    ("TOTAL_OUTSTANDING", "Total sisa pokok pinjaman", "SUM(outstanding)"),
    ("JUMLAH_NASABAH", "Jumlah nasabah aktif", "COUNT(nasabah_id) WHERE aktif"),
    ("VOL_TRANSAKSI", "Total volume transaksi", "SUM(jumlah) WHERE berhasil"),
]
y = Inches(1.6)
for name, desc, formula in metrics:
    add_text_box(slide, Inches(1), y, Inches(3), Inches(0.6), name, 16, True, BLUE)
    add_text_box(slide, Inches(4.2), y, Inches(5.5), Inches(0.6), desc, 15, False, DARK)
    add_text_box(slide, Inches(10), y, Inches(5.5), Inches(0.6), formula, 13, False, GRAY)
    y += Inches(0.7)

add_text_box(slide, Inches(1), Inches(7), Inches(14), Inches(1),
             "Semua metrics menggunakan terminologi OJK (Otoritas Jasa Keuangan)\n"
             "Kolektibilitas: 1-Lancar, 2-DPK, 3-Kurang Lancar, 4-Diragukan, 5-Macet",
             16, False, ORANGE, PP_ALIGN.CENTER)

# ========== SLIDE 11: Part 4 - Cortex Agent ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(16), Inches(1.2), DARK_BLUE)
add_text_box(slide, Inches(1), Inches(0.2), Inches(14), Inches(0.8), "Part 4: Cortex Agent - AI Banking Assistant", 32, True, WHITE)

add_text_box(slide, Inches(1), Inches(1.5), Inches(14), Inches(0.6), "4 Tools dalam 1 Agent:", 22, True, PURPLE)

tools = [
    ("Cortex Search", "Menjawab dari dokumen SOP bank", "SOP, KYC, AML, Kredit, Digital"),
    ("Cortex Analyst", "Query data operasional via NL", "NPL, DPK, transaksi, nasabah"),
    ("Web Search", "Informasi terkini dari internet", "Regulasi OJK, suku bunga BI"),
    ("Send Email", "Kirim notifikasi & laporan", "Alert NPL, laporan harian"),
]

y = Inches(2.4)
colors = [GREEN, BLUE, ORANGE, RED]
for i, (name, desc, examples) in enumerate(tools):
    shape = add_shape(slide, Inches(1), y, Inches(14), Inches(1.1), RGBColor(0xF9, 0xFA, 0xFB))
    shape.line.color.rgb = colors[i]
    shape.line.width = Pt(2)
    add_text_box(slide, Inches(1.3), y + Inches(0.1), Inches(3), Inches(0.9), name, 18, True, colors[i])
    add_text_box(slide, Inches(4.5), y + Inches(0.1), Inches(5), Inches(0.9), desc, 16, False, DARK)
    add_text_box(slide, Inches(10), y + Inches(0.1), Inches(4.5), Inches(0.9), examples, 14, False, GRAY)
    y += Inches(1.3)

add_text_box(slide, Inches(1), Inches(7.5), Inches(14), Inches(0.6),
             "Agent secara otomatis memilih tool yang tepat berdasarkan pertanyaan user",
             17, True, PURPLE, PP_ALIGN.CENTER)

# ========== SLIDE 12: Answer Types ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(16), Inches(1.2), DARK_BLUE)
add_text_box(slide, Inches(1), Inches(0.2), Inches(14), Inches(0.8), "Jenis Jawaban Agent", 32, True, WHITE)

answer_types = [
    ("Fact-Based Answer", '"Total DPK saat ini Rp 2.5 Triliun, naik 3% dari bulan lalu"', BLUE),
    ("Diagnostic Answer", '"NPL ratio naik 2% karena sektor F&B dan Properti terdampak penurunan ekonomi"', ORANGE),
    ("Recommendation", '"Rekomendasi: fokus restrukturisasi di 3 cabang Jawa Barat dengan NPL > 8%"', GREEN),
    ("Chart / Visualization", 'Tabel, bar chart, line chart untuk memvisualisasikan data', PURPLE),
    ("Email Notification", '"Laporan NPL telah dikirim ke risk.manager@banknusantara.co.id"', RED),
]

y = Inches(1.6)
for title, example, color in answer_types:
    shape = add_shape(slide, Inches(1), y, Inches(14), Inches(1.2), RGBColor(0xF9, 0xFA, 0xFB))
    shape.line.color.rgb = color
    shape.line.width = Pt(2)
    add_text_box(slide, Inches(1.3), y + Inches(0.05), Inches(4), Inches(0.5), title, 18, True, color)
    add_text_box(slide, Inches(1.3), y + Inches(0.55), Inches(13.4), Inches(0.5), example, 15, False, GRAY)
    y += Inches(1.35)

# ========== SLIDE 13: Demo Scenarios ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(16), Inches(1.2), DARK_BLUE)
add_text_box(slide, Inches(1), Inches(0.2), Inches(14), Inches(0.8), "Demo: Contoh Pertanyaan ke Agent", 32, True, WHITE)

demos = [
    ("Unstructured", '"Apa syarat pembukaan rekening korporasi?"', "-> Cortex Search"),
    ("Structured", '"Berapa NPL ratio per jenis kredit?"', "-> Cortex Analyst + Chart"),
    ("Diagnostic", '"Sektor ekonomi mana yang NPL-nya paling tinggi dan kenapa?"', "-> Analyst + AI Analysis"),
    ("Web Search", '"Berapa suku bunga acuan BI saat ini?"', "-> Web Search"),
    ("Combined", '"Jelaskan prosedur restrukturisasi dan berapa kredit yang perlu direstrukturisasi?"', "-> Search + Analyst"),
    ("Email", '"Kirim ringkasan NPL ke risk.manager@bank.co.id"', "-> Analyst + Email"),
]

y = Inches(1.5)
for label, question, tool in demos:
    add_text_box(slide, Inches(1), y, Inches(2.5), Inches(0.5), label, 16, True, BLUE)
    add_text_box(slide, Inches(3.5), y, Inches(8.5), Inches(0.5), question, 16, False, DARK)
    add_text_box(slide, Inches(12.5), y, Inches(3), Inches(0.5), tool, 14, False, GREEN)
    y += Inches(0.65)

add_text_box(slide, Inches(1), Inches(6.5), Inches(14), Inches(1.5),
             "AHA Moment:\n"
             "Dalam 1 session, kita membangun AI Assistant yang menggabungkan\n"
             "knowledge base dokumen + data warehouse + internet - semua di dalam Snowflake!",
             20, True, PURPLE, PP_ALIGN.CENTER)

# ========== SLIDE 14: Part 5 - Snowflake Intelligence ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(16), Inches(1.2), DARK_BLUE)
add_text_box(slide, Inches(1), Inches(0.2), Inches(14), Inches(0.8), "Part 5: Snowflake Intelligence", 32, True, WHITE)

add_text_box(slide, Inches(1), Inches(1.5), Inches(14), Inches(0.6),
             "Publish Agent sebagai Aplikasi untuk Business Users", 22, True, TEAL)

users = [
    ("Relationship Manager", "Ringkasan nasabah, SOP kredit,\ndata portofolio untuk meeting", Inches(1)),
    ("Risk Analyst", "NPL monitoring, early warning,\nanalisis sektor ekonomi", Inches(5.7)),
    ("Branch Manager", "Performa cabang, target vs aktual,\nranking antar cabang", Inches(10.4)),
]

for name, desc, left in users:
    shape = add_shape(slide, left, Inches(2.5), Inches(4.3), Inches(2.5), RGBColor(0xF0, 0xFD, 0xF4))
    shape.line.color.rgb = TEAL
    shape.line.width = Pt(2)
    add_text_box(slide, left + Inches(0.3), Inches(2.7), Inches(3.7), Inches(0.5), name, 18, True, TEAL, PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.3), Inches(3.4), Inches(3.7), Inches(1.2), desc, 15, False, DARK, PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(5.5), Inches(14), Inches(2),
             "Key Benefits:\n\n"
             "  - No SQL knowledge required - cukup tanya dalam bahasa natural\n"
             "  - Role-based access - setiap user hanya lihat data yang boleh diakses\n"
             "  - Multi-source answers - SOP + Database + Internet dalam satu jawaban\n"
             "  - Governance built-in - semua data tetap di dalam Snowflake",
             17, False, DARK)

# ========== SLIDE 15: Summary ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(16), Inches(1.2), DARK_BLUE)
add_text_box(slide, Inches(1), Inches(0.2), Inches(14), Inches(0.8), "Ringkasan: Apa yang Sudah Kita Bangun", 32, True, WHITE)

summary = [
    ("PARSE_DOCUMENT", "Extract text dari 6 PDF dokumen SOP bank"),
    ("SPLIT_TEXT", "Chunking dokumen untuk pencarian optimal"),
    ("Cortex Search", "Hybrid search pada dokumen SOP internal"),
    ("Semantic View", "Data model 5 tabel dengan 8 metrics perbankan"),
    ("Cortex Analyst", "Natural language query ke data operasional"),
    ("Cortex Agent", "Multi-tool AI assistant (Search + Analyst + Web + Email)"),
    ("Snowflake Intelligence", "Published agent untuk business users"),
]

y = Inches(1.5)
for feature, desc in summary:
    add_text_box(slide, Inches(1.5), y, Inches(4), Inches(0.5), feature, 17, True, BLUE)
    add_text_box(slide, Inches(5.5), y, Inches(9.5), Inches(0.5), desc, 17, False, DARK)
    y += Inches(0.6)

add_text_box(slide, Inches(1), Inches(7.0), Inches(14), Inches(1),
             "Semua berjalan di dalam Snowflake - Zero Infrastructure Management",
             22, True, GREEN, PP_ALIGN.CENTER)

# ========== SLIDE 16: Thank You ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BLUE)
add_shape(slide, Inches(0), Inches(0), Inches(16), Inches(0.15), SNOW_BLUE)
add_shape(slide, Inches(0), Inches(8.85), Inches(16), Inches(0.15), SNOW_BLUE)

add_text_box(slide, Inches(2), Inches(2.5), Inches(12), Inches(1.5),
             "Hands-On Lab Time!", 48, True, WHITE, PP_ALIGN.CENTER)
add_text_box(slide, Inches(2), Inches(4.2), Inches(12), Inches(1),
             "Ikuti step-by-step di README.md", 24, False, SNOW_BLUE, PP_ALIGN.CENTER)
add_text_box(slide, Inches(2), Inches(5.5), Inches(12), Inches(1),
             "Jangan ragu untuk bertanya!", 22, False, RGBColor(0xBF, 0xDB, 0xFE), PP_ALIGN.CENTER)
add_text_box(slide, Inches(2), Inches(7.0), Inches(12), Inches(0.5),
             "github.com/arzamuhammad/snowflake_bank_workshop", 16, False, GRAY, PP_ALIGN.CENTER)

out_path = os.path.join(os.path.dirname(__file__), "slides", "Day1_Session1_RAG_Snowflake_Intelligence.pptx")
prs.save(out_path)
print(f"Saved: {out_path}")
